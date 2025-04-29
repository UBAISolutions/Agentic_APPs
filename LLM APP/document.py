import pandas as pd
import PyPDF2
import docx
from pptx import Presentation
import os
import re
from typing import List, Dict, Union, Any

class DocumentProcessor:
    def __init__(self):
        self.supported_extensions = ['csv', 'xlsx', 'pdf', 'pptx', 'docx']
    
    def read_file(self, file_path: str) -> Dict[str, Any]:
        """Read a file and extract its content based on file type."""
        _, extension = os.path.splitext(file_path)
        extension = extension[1:].lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file format: {extension}")
        
        content = {}
        
        if extension == 'csv':
            content['data'] = pd.read_csv(file_path)
            content['text'] = self._dataframe_to_text(content['data'])
            content['type'] = 'tabular'
        
        elif extension == 'xlsx':
            content['data'] = pd.read_excel(file_path, sheet_name=None)
            content['text'] = ""
            for sheet_name, sheet_data in content['data'].items():
                content['text'] += f"Sheet: {sheet_name}\n"
                content['text'] += self._dataframe_to_text(sheet_data)
            content['type'] = 'tabular'
        
        elif extension == 'pdf':
            content['text'] = self._extract_text_from_pdf(file_path)
            content['type'] = 'document'
        
        elif extension == 'docx':
            content['text'] = self._extract_text_from_docx(file_path)
            content['type'] = 'document'
        
        elif extension == 'pptx':
            content['text'] = self._extract_text_from_pptx(file_path)
            content['type'] = 'presentation'
            
        return content
    
    def _dataframe_to_text(self, df: pd.DataFrame) -> str:
        """Convert DataFrame to text representation."""
        text = f"Columns: {', '.join(df.columns.tolist())}\n"
        text += f"Data shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
        text += "Sample data:\n"
        text += df.head(5).to_string() + "\n\n"
        return text
    
    def _extract_text_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        with open(file_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page_num in range(len(reader.pages)):
                text += reader.pages[page_num].extract_text() + "\n"
        return text
    
    def _extract_text_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    
    def _extract_text_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        presentation = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(presentation.slides):
            text += f"Slide {slide_num + 1}:\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
            text += "\n"
        return text
    
    def chunk_text(self, text: str, chunk_size: int = 1000, overlap: int = 100) -> List[str]:
        """Split text into overlapping chunks."""
        chunks = []
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            if chunk:
                chunks.append(chunk)
        return chunks
    
    def process_document(self, file_path: str, chunk_size: int = 1000, overlap: int = 100) -> Dict[str, Any]:
        """Process a document: extract content and chunk it."""
        content = self.read_file(file_path)
        
        if content['type'] in ['document', 'presentation']:
            content['chunks'] = self.chunk_text(content['text'], chunk_size, overlap)
        else:  # tabular data
            # Keep the original data for analysis
            content['chunks'] = [content['text']]
            
        return content
